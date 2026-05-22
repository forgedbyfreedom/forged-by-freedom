#!/usr/bin/env python3
"""Generate the Acoustic Drone Detection capabilities deck (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme ────────────────────────────────────────────────────────────
BG     = RGBColor(0x0E, 0x12, 0x18)   # near-black slate
PANEL  = RGBColor(0x18, 0x20, 0x2A)
ACCENT = RGBColor(0xFF, 0x7A, 0x18)   # orange
ACCENT2= RGBColor(0x35, 0xC2, 0x8E)   # teal/green
TEXT   = RGBColor(0xE9, 0xED, 0xF1)
MUTE   = RGBColor(0x9A, 0xA6, 0xB2)

W, H = Inches(13.333), Inches(7.5)     # 16:9

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H)  # rectangle
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None):
    shp = s.shapes.add_shape(1, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor="t"):
    """runs: list of (string, size, color, bold) or list-of-such for paragraphs."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"t": 1, "m": 3, "b": 4}[anchor]
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    return tb


def bullets(s, x, y, w, h, items, size=16, gap=8):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        lvl, txt, color, bold = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        bullet = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet + txt
        r.font.size = Pt(size - lvl * 2); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    box(s, 0, 0, W, Inches(1.15), PANEL)
    box(s, 0, Inches(1.15), W, Pt(3), ACCENT)
    text(s, Inches(0.5), Inches(0.12), Inches(12), Inches(0.4),
         [(kicker, 12, ACCENT, True)])
    text(s, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.7),
         [(title, 28, TEXT, True)])


# ── 1. Title ─────────────────────────────────────────────────────────
s = slide()
box(s, 0, Inches(2.5), W, Pt(3), ACCENT)
text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.6),
     [(("Acoustic Drone Detection System"), 44, TEXT, True)])
text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.6),
     [("Passive multi-layer detection • RF-fusion ready • laptop proof-of-concept",
       20, ACCENT2, False)])
text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
     [("Capabilities, Range & Phased Roadmap", 16, MUTE, False)])
text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
     [("Prepared for Forged by Freedom  •  Prototype / planning document",
       12, MUTE, False)])

# ── 2. Executive summary ─────────────────────────────────────────────
s = slide(); header(s, "OVERVIEW", "Executive Summary")
bullets(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6), [
    (0, "Goal: passively detect drones by their acoustic signature, with no transmissions of our own.", TEXT, True),
    (0, "Core principle: a drone is not one tone — it is a structured, persistent, modulated harmonic stack. We detect the structure, not just sound in a band.", TEXT, False),
    (0, "Five cooperating layers reject the everyday confusers (fans, leaf blowers, RC planes, HVAC, traffic, insects):", ACCENT2, True),
    (1, "Harmonic stack — blade-pass fundamental + its overtones", TEXT, False),
    (1, "Persistence — must hold for several seconds", TEXT, False),
    (1, "Multi-rotor beat — multiple props slightly out of sync", TEXT, False),
    (1, "RPM modulation — the tell-tale hover 'wobble'", TEXT, False),
    (1, "Direction-of-arrival — bearing from a mic array", TEXT, False),
    (0, "Designed to fuse with a separate RF detection system: RF gives long-range early warning, acoustic confirms and catches RF-silent / autonomous drones.", TEXT, False),
    (0, "Working software prototype already validated on synthetic + recorded audio. Laptop proof-of-concept needs ~$40–110 in parts.", ACCENT, True),
])

# ── 3. The problem ───────────────────────────────────────────────────
s = slide(); header(s, "WHY IT'S HARD", "Frequency Alone Always False-Alarms")
bullets(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.6), [
    (0, "Propeller noise overlaps heavily with:", TEXT, True),
    (1, "Fans, HVAC, compressors", MUTE, False),
    (1, "Leaf blowers, lawn equipment", MUTE, False),
    (1, "RC planes / other hobby craft", MUTE, False),
    (1, "Traffic, insects, wind", MUTE, False),
    (0, "Monitoring 'one clean band' cannot isolate drones — there is no clean band.", ACCENT, True),
])
box(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), PANEL)
bullets(s, Inches(7.2), Inches(1.75), Inches(5.2), Inches(5.0), [
    (0, "The fix: require a PATTERN, not a frequency", ACCENT2, True),
    (1, "steady multi-tone harmonic stack", TEXT, False),
    (1, "repeated signature lasting seconds", TEXT, False),
    (1, "rapid RPM modulation", TEXT, False),
    (1, "multiple props slightly out of sync", TEXT, False),
    (1, "directional confirmation", TEXT, False),
    (1, "cross-check with RF / video", TEXT, False),
    (0, "Each layer a confuser must defeat simultaneously — which they almost never do.", TEXT, False),
])

# ── 4. Physics ───────────────────────────────────────────────────────
s = slide(); header(s, "THE SIGNATURE", "Acoustic Physics — Blade-Pass Frequency")
bullets(s, Inches(0.6), Inches(1.45), Inches(6.1), Inches(5.7), [
    (0, "Blade-pass frequency (BPF) = the fundamental tone:", TEXT, True),
    (1, "BPF = RPM ÷ 60 × number of blades", ACCENT, True),
    (0, "Examples:", TEXT, True),
    (1, "2-blade @ 6,000 RPM = 200 Hz", MUTE, False),
    (1, "3-blade @ 6,000 RPM = 300 Hz", MUTE, False),
    (1, "Big slow heavy-lift prop ≈ 60–250 Hz", MUTE, False),
    (0, "Then the harmonics: 2×, 3×, 4×, 5× BPF, often reaching 1–5 kHz.", TEXT, False),
    (0, "Motor/ESC whine: ~4–12 kHz (close-in only — dies fast outdoors).", TEXT, False),
])
box(s, Inches(6.95), Inches(1.45), Inches(5.75), Inches(5.3), PANEL)
bullets(s, Inches(7.2), Inches(1.7), Inches(5.3), Inches(5.0), [
    (0, "Band plan", ACCENT2, True),
    (1, "Low  50–700 Hz — blade-pass + low harmonics (range workhorse)", TEXT, False),
    (1, "Mid  700 Hz–5 kHz — harmonic stack", TEXT, False),
    (1, "High 5–12 kHz — ESC/motor whine (short range)", TEXT, False),
    (0, "Low frequencies barely attenuate in air → they carry far. High frequencies fade within ~100 m outdoors.", TEXT, False),
    (0, "Heavy-lift = low fundamental + loud = excellent acoustic target.", ACCENT, True),
])

# ── 5. Five layers ───────────────────────────────────────────────────
s = slide(); header(s, "HOW IT WORKS", "The Five Detection Layers")
cards = [
    ("1 · Harmonic stack", "Find a fundamental in 50–700 Hz with ≥3 overtones above the local noise floor.", ACCENT),
    ("2 · Persistence", "Signature must reappear for ~3 s before alerting. Kills one-off transients.", ACCENT2),
    ("3 · Multi-rotor beat", "Two fundamentals a few Hz apart = multiple desynced rotors (hexa/octo).", ACCENT),
    ("4 · RPM modulation", "Hover causes a small rapid 'wobble' in the fundamental. Steady tones do not.", ACCENT2),
    ("5 · Direction (DOA)", "Mic array + SRP-PHAT gives a bearing and rejects diffuse background noise.", ACCENT),
    ("Fusion-ready", "Outputs structured detections for cross-checking with the RF system.", MUTE),
]
x0, y0, cw, ch, gx, gy = Inches(0.55), Inches(1.5), Inches(4.0), Inches(2.45), Inches(0.2), Inches(0.25)
for i, (t, d, c) in enumerate(cards):
    col, row = i % 3, i // 3
    cx = x0 + col * (cw + gx); cy = y0 + row * (ch + gy)
    box(s, cx, cy, cw, ch, PANEL)
    box(s, cx, cy, cw, Pt(4), c)
    text(s, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), Inches(0.5),
         [(t, 16, c, True)])
    text(s, cx + Inches(0.2), cy + Inches(0.75), cw - Inches(0.4), Inches(1.5),
         [(d, 13, TEXT, False)])

# ── 6. POC architecture ──────────────────────────────────────────────
s = slide(); header(s, "PROOF OF CONCEPT", "Laptop Setup — Signal Chain")
chain = ["USB mic\n(or pair)", "Laptop\naudio in", "FFT +\nharmonic\nanalysis",
         "5-layer\nlogic", "Detection\n+ bearing", "→ RF system\nfusion"]
n = len(chain); cw = Inches(1.85); gap = Inches(0.18)
total = n * cw + (n - 1) * gap
x = (W - total) / 2; y = Inches(2.6)
for i, label in enumerate(chain):
    c = ACCENT if i in (2, 3) else PANEL
    b = box(s, x, y, cw, Inches(1.5), PANEL)
    box(s, x, y, cw, Pt(4), ACCENT if i in (2, 3, 4) else ACCENT2)
    text(s, x, y + Inches(0.2), cw, Inches(1.2),
         [(label, 13, TEXT, True)], align=PP_ALIGN.CENTER, anchor="m")
    if i < n - 1:
        text(s, x + cw, y + Inches(0.45), gap, Inches(0.5),
             [("→", 20, ACCENT, True)], align=PP_ALIGN.CENTER)
    x += cw + gap
bullets(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.4), [
    (0, "Pure-numpy analysis path runs on any modern laptop — no GPU, no cloud, no API keys.", TEXT, False),
    (0, "Validate today for $0: feed it a recorded drone flyby (YouTube → 16-bit WAV) before buying any hardware.", ACCENT2, True),
    (0, "Add one USB mic for live single-channel detection; add a 4-mic array for live bearing.", TEXT, False),
])

# ── 7. POC parts list ────────────────────────────────────────────────
s = slide(); header(s, "BILL OF MATERIALS", "Proof-of-Concept on a Laptop")
rows = [
    ("Item", "Purpose", "Cost (USD)", True),
    ("Laptop you already own", "Runs the detector (numpy). No GPU needed.", "$0", False),
    ("Recorded drone clips → WAV", "Validate the DSP before buying anything.", "$0", False),
    ("Dayton iMM-6 (or USB MEMS mic)", "Single-channel live measurement mic.", "$25–60", False),
    ("USB mic OR ReSpeaker 4-Mic Array", "4 mics = live bearing (direction).", "$70 (array)", False),
    ("Furry windscreen / dead-cat", "Mandatory outdoors — kills wind in low band.", "$10–25", False),
    ("USB extension / tripod clip", "Place mic away from laptop fan noise.", "$10", False),
    ("TOTAL — single mic POC", "Prove detection + RPM/beat layers", "≈ $40–95", True),
    ("TOTAL — array POC (with DOA)", "Adds live direction-of-arrival", "≈ $90–110", True),
]
ty = Inches(1.55); rh = Inches(0.55)
cols = [Inches(0.6), Inches(4.7), Inches(9.7)]
cwid = [Inches(4.1), Inches(5.0), Inches(3.0)]
for i, (a, b, c, bold) in enumerate(rows):
    yy = ty + i * rh
    fill = PANEL if (i == 0 or bold) else None
    if fill: box(s, Inches(0.5), yy, Inches(12.3), rh, fill)
    col = ACCENT if (i == 0 or bold) else TEXT
    for cx, cwd, val, al in [(cols[0], cwid[0], a, PP_ALIGN.LEFT),
                             (cols[1], cwid[1], b, PP_ALIGN.LEFT),
                             (cols[2], cwid[2], c, PP_ALIGN.RIGHT)]:
        text(s, cx, yy + Inches(0.05), cwd, rh,
             [(val, 13, col, (i == 0 or bold))], align=al, anchor="m")

# ── 8. Range — acoustic ──────────────────────────────────────────────
s = slide(); header(s, "PERFORMANCE", "Detection Range — Acoustic (night / quiet)")
bullets(s, Inches(0.6), Inches(1.45), Inches(6.1), Inches(5.7), [
    (0, "Small consumer drone (~70 dB @1m):", TEXT, True),
    (1, "Quiet rural night: ~150–300 m", MUTE, False),
    (1, "Quiet suburban: ~80–150 m", MUTE, False),
    (0, "Heavy-lift electric + payload (~85–95 dB @1m):", ACCENT, True),
    (1, "Quiet rural night: ~0.8–1.5 km", TEXT, False),
    (1, "Loudest builds, dead calm: ~2 km", TEXT, False),
    (1, "Quiet suburban: ~300–700 m", TEXT, False),
    (0, "Classification range ≈ 60–70% of raw detection range.", MUTE, False),
])
box(s, Inches(6.95), Inches(1.45), Inches(5.75), Inches(5.3), PANEL)
bullets(s, Inches(7.2), Inches(1.7), Inches(5.3), Inches(5.0), [
    (0, "Range multipliers", ACCENT2, True),
    (1, "Parabolic / shotgun mic: +10–20 dB in-beam → 2–3× range in that sector", TEXT, False),
    (1, "4-mic array beamforming: ~6 dB + noise rejection", TEXT, False),
    (1, "Night + quiet = lowest noise floor = best case", TEXT, False),
    (0, "Range killers", ACCENT, True),
    (1, "Wind (raises low-band floor — use windscreen)", TEXT, False),
    (1, "Terrain / obstructions", TEXT, False),
])

# ── 9. Range — fusion ────────────────────────────────────────────────
s = slide(); header(s, "PERFORMANCE", "Combined with RF — Fusion Envelope")
rows = [
    ("Target", "Effective range (night/quiet)", "Driven by", True),
    ("RF-emitting consumer drone", "~1–3 km early warning; acoustic confirms last 150–300 m", "RF", False),
    ("Heavy-lift, RF-emitting", "RF 1–5 km + acoustic 0.8–2 km (heavy overlap)", "RF + acoustic", False),
    ("RF-silent / autonomous drone", "~150 m–2 km depending on size", "Acoustic only", False),
]
ty = Inches(1.7); rh = Inches(0.85)
cols = [Inches(0.6), Inches(4.4), Inches(10.4)]; cwid = [Inches(3.7), Inches(5.8), Inches(2.3)]
for i, (a, b, c, bold) in enumerate(rows):
    yy = ty + i * rh
    box(s, Inches(0.5), yy, Inches(12.3), rh, PANEL if i == 0 else (None if i % 2 else PANEL))
    col = ACCENT if i == 0 else TEXT
    for cx, cwd, val, al in [(cols[0], cwid[0], a, PP_ALIGN.LEFT),
                             (cols[1], cwid[1], b, PP_ALIGN.LEFT),
                             (cols[2], cwid[2], c, PP_ALIGN.LEFT)]:
        text(s, cx, yy, cwd, rh, [(val, 14, col, i == 0)], align=al, anchor="m")
text(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.8),
     [[("Fusion's payoff is confidence + coverage, not just range: ", 15, ACCENT2, True),
       ("a track confirmed on both modalities is high-trust and low-false-alarm, and the union means you are blind to neither the RF-silent drone (acoustic catches it) nor the distant emitter (RF catches it).", 15, TEXT, False)]])

# ── 10. Limitations ──────────────────────────────────────────────────
s = slide(); header(s, "HONEST LIMITS", "What This Does Not Do (yet)")
bullets(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6), [
    (0, "Acoustic is a close-in layer, not perimeter radar — range is tens of meters to ~2 km, weather-dependent.", TEXT, False),
    (0, "Bearing from a small array is coarse at low frequency (physics, not code). A larger baseline or high-harmonic steering is needed for tight bearings.", TEXT, False),
    (0, "Multi-rotor 'beat' flag is reliable as binary; exact beat value is coarse (FFT resolution).", TEXT, False),
    (0, "Wind degrades performance fast — windscreen + siting matter as much as the mic.", TEXT, False),
    (0, "Acoustic alone will still occasionally confuse loud RC craft / equipment — that's why RF cross-check matters before acting.", TEXT, False),
    (0, "No classification of make/model yet (that's a future ML phase). Current output: 'drone-like acoustic signature, bearing X, multi-rotor yes/no.'", ACCENT, True),
])

# ── 11. Roadmap ──────────────────────────────────────────────────────
s = slide(); header(s, "ROADMAP", "Phased Development Plan")
phases = [
    ("Phase 0", "Laptop POC", "Validate DSP on recordings + 1 mic. $0–95.", ACCENT2),
    ("Phase 1", "Single field node", "Raspberry Pi + 4-mic array, weatherproof. Live bearing. ~$180–250.", ACCENT),
    ("Phase 2", "Directional / range node", "Parabolic or shotgun mic for 2 km+ in a sector. ~$350–650.", ACCENT2),
    ("Phase 3", "Multi-node network", "3+ nodes triangulate position, not just bearing. Mesh / LoRa backhaul.", ACCENT),
    ("Phase 4", "RF + acoustic fusion", "Unified track table; correlate RF DF with acoustic bearing.", ACCENT2),
    ("Phase 5", "ML classification + cueing", "Train classifier on signatures; auto-slew a PTZ camera to bearing.", ACCENT),
]
y = Inches(1.55); rh = Inches(0.86)
for i, (p, t, d, c) in enumerate(phases):
    yy = y + i * rh
    box(s, Inches(0.5), yy, Inches(12.3), Inches(0.78), PANEL)
    box(s, Inches(0.5), yy, Pt(5), Inches(0.78), c)
    text(s, Inches(0.75), yy, Inches(1.6), Inches(0.78), [(p, 15, c, True)], anchor="m")
    text(s, Inches(2.4), yy, Inches(2.9), Inches(0.78), [(t, 15, TEXT, True)], anchor="m")
    text(s, Inches(5.3), yy, Inches(7.3), Inches(0.78), [(d, 13, MUTE, False)], anchor="m")

# ── 12. Future expansion ─────────────────────────────────────────────
s = slide(); header(s, "FUTURE", "Expansion Opportunities")
bullets(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.6), [
    (0, "Sensing", ACCENT, True),
    (1, "ML signature classifier (make/model/payload-state)", TEXT, False),
    (1, "High-harmonic DOA for tight bearings", TEXT, False),
    (1, "Optical / IR camera slewed to acoustic bearing", TEXT, False),
    (1, "Rotor-count + RPM estimation from beat structure", TEXT, False),
])
bullets(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.6), [
    (0, "System", ACCENT2, True),
    (1, "Multi-node TDOA → true GPS position + track", TEXT, False),
    (1, "Unified fusion dashboard with RF DF", TEXT, False),
    (1, "Mesh / LoRa backhaul for remote nodes", TEXT, False),
    (1, "Alerting: SMS / app / siren / log", TEXT, False),
    (1, "Solar + battery for off-grid perimeter nodes", TEXT, False),
])

# ── 13. Cost summary ─────────────────────────────────────────────────
s = slide(); header(s, "BUDGET", "Indicative Cost by Phase")
rows = [
    ("Phase", "What you get", "Indicative cost", True),
    ("0 · Laptop POC", "Prove the detection works", "$0–95", False),
    ("1 · Single node", "24/7 field node, live bearing", "$180–250", False),
    ("2 · Directional node", "Long-range (2 km+) sector coverage", "$350–650", False),
    ("3 · 3-node network", "Position/triangulation", "~$700–1,500", False),
    ("4 · RF fusion", "Software integration (uses existing RF kit)", "Mostly dev time", False),
    ("5 · ML + cueing", "Classification + camera slew", "Dev + ~$150–600 optics", False),
]
ty = Inches(1.7); rh = Inches(0.72)
cols = [Inches(0.6), Inches(4.6), Inches(10.0)]; cwid = [Inches(3.9), Inches(5.3), Inches(2.7)]
for i, (a, b, c, bold) in enumerate(rows):
    yy = ty + i * rh
    box(s, Inches(0.5), yy, Inches(12.3), rh, PANEL if (i == 0 or i % 2 == 0) else None)
    col = ACCENT if i == 0 else TEXT
    for cx, cwd, val, al in [(cols[0], cwid[0], a, PP_ALIGN.LEFT),
                             (cols[1], cwid[1], b, PP_ALIGN.LEFT),
                             (cols[2], cwid[2], c, PP_ALIGN.RIGHT)]:
        text(s, cx, yy, cwd, rh, [(val, 13, col, i == 0)], align=al, anchor="m")
text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
     [("Costs are indicative hobby/prototype pricing, USD; exclude labor. Confirm current part prices before purchase.", 11, MUTE, False)])

# ── 14. Next steps ───────────────────────────────────────────────────
s = slide(); header(s, "ACTION", "Recommended Next Steps")
bullets(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.4), [
    (0, "1. Validate for $0 — run the detector against 3–5 recorded heavy-lift drone clips. Confirm it locks the harmonic stack + flags multi-rotor.", TEXT, True),
    (0, "2. Buy the single-mic POC kit (~$40–95) and test live against a known drone at measured distances to calibrate real-world range.", TEXT, False),
    (0, "3. Add the 4-mic array to bring up live bearing; tune band weighting for your specific heavy-lift fleet.", TEXT, False),
    (0, "4. Define the fusion interface to the RF system (shared track format, time-sync).", TEXT, False),
    (0, "5. Decide Phase 1 siting — where the field node lives, power, and backhaul.", TEXT, False),
    (0, "Software prototype (drone_detect.py) is attached and already tested.", ACCENT, True),
])

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "Acoustic_Drone_Detection_System.pptx"
prs.save(out)
print("saved", out, "—", len(prs.slides._sldIdLst), "slides")
